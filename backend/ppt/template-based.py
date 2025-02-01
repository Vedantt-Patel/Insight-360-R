import re
import os
# from dotenv import load_dotenv
import random
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Pt, Inches
import sys
# load_dotenv()
template_choice = sys.argv[3] 
print(template_choice)
def parse_llm_output(llm_response):
    slides = re.split(r'#Slide:\s*\d+', llm_response.strip())  
    slide_numbers = re.findall(r'#Slide:\s*(\d+)', llm_response)  
    parsed_slides = []

    for idx, slide in enumerate(slides[1:]):  
        lines = slide.strip().split('\n')
        
        header_line = next((line for line in lines if line.startswith("#Header:")), None)
        header = header_line.split(": ", 1)[1] if header_line else "No Header"

        image_line = next((line for line in lines if line.startswith("#Image:")), None)
        image = image_line.split(": ", 1)[1].strip().lower() == "true" if image_line else False

        content_lines = [line for line in lines if line.startswith("#Content:") or not line.startswith("#")]
        content = "\n".join(line.replace("#Content:", "").strip() for line in content_lines).strip()

        parsed_slides.append({
            "slide_number": int(slide_numbers[idx]),
            "title": header,
            "content": content,
            "image": image
        })

    return parsed_slides

with open("generated_slides/final_presentation_slides.txt", "r", encoding="utf-8") as file:
    content = file.read()

slides_content = parse_llm_output(content)
    
with open("generated_slides/title_authors.txt", "r", encoding="utf-8") as f:
    for line in f:
        line = line.strip()
        if line.startswith("Title:"):
            title = line.replace("Title:", "").strip()
        elif line.startswith("Authors:"):
            authors = line.replace("Authors:", "").strip()

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import os

def apply_template_styles(placeholder, template_choice, is_title=True):
    templates = {
        'simple': {'title_font': 'Calibri Light (Headings)', 'subtitle_font': 'Calibri (Body)',
                       'title_size': (50, 36), 'subtitle_size': (26, 18), 'color': RGBColor(0, 0, 0)},
        'futuristic': {'title_font': 'Calibri Light (Headings)', 'subtitle_font': 'Calibri (Body)',
                       'title_size': (40, 36), 'subtitle_size': (18, 18), 'color': RGBColor(255, 255, 255)},
        'dark': {'title_font': 'Walbaum Display', 'subtitle_font': 'Gill Sans MT',
                 'title_size': (44, 18), 'subtitle_size': (18, 18), 'color': RGBColor(255, 255, 255)},
        'minimal': {'title_font': 'Meiryo', 'subtitle_font': 'Meiryo',
                    'title_size': (28, 14), 'subtitle_size': (14, 14), 'color': RGBColor(0, 0, 0)},
        'modern': {'title_font': 'Tw Cen MT', 'subtitle_font': 'Segoe UI Light',
                   'title_size': (28, 14), 'subtitle_size': (14, 14), 'color': RGBColor(255, 255, 255)},
        'dark_modern': {'title_font': 'Britannic Bold', 'subtitle_font': 'Sylfaen',
                        'title_size': (40, 20), 'subtitle_size': (20, 20), 'color': RGBColor(255, 255, 255)},
        'bright_modern': {'title_font': 'Tw Cen MT', 'subtitle_font': 'Tw Cen MT',
                          'title_size': (40, 20), 'subtitle_size': (18, 18), 'color': RGBColor(0, 0, 0)},
        'geometric': {'title_font': 'Tw Cen MT', 'subtitle_font': 'Tw Cen MT',
                      'title_size': (60, 54), 'subtitle_size': (24, 24), 'color': RGBColor(0, 0, 0)},
    }
    
    template = templates.get(template_choice, templates['modern'])
    font_name = template['title_font'] if is_title else template['subtitle_font']
    font_size = template['title_size'][0] if is_title else template['subtitle_size'][0]
    font_color = template['color']
    
    for paragraph in placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color

def adjust_placeholder_size(placeholder, left, top, width, height):
    placeholder.left = Inches(left)
    placeholder.top = Inches(top)
    placeholder.width = Inches(width)
    placeholder.height = Inches(height)

def get_random_image(image_folder):
    if os.path.exists(image_folder):
        images = [f for f in os.listdir(image_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
        return os.path.join(image_folder, random.choice(images)) if images else None
    return None

def create_ppt(slides_content, template_choice, presentation_title, presenter_name, image_folder='images'):
    c=0
    template_path = os.path.join('templates', f"{template_choice}.pptx")
    prs = Presentation(template_path)
    
    layouts = {
        'left_image': {
            'title': {'left': 0.5, 'top': 0.3, 'width': 12, 'height': 0.8},
            'content': {'left': 5.8, 'top': 1.3, 'width': 7, 'height': 5.7},
            'image': {'left': 0.5, 'top': 1.3, 'width': 4.8, 'height': 5.7}
        },
        'right_image': {
            'title': {'left': 0.5, 'top': 0.3, 'width': 12, 'height': 0.8},
            'content': {'left': 0.5, 'top': 1.3, 'width': 7, 'height': 5.7},
            'image': {'left': 8, 'top': 1.3, 'width': 4.8, 'height': 5.7}
        },
        'no_image': {
            'title': {'left': 0.5, 'top': 0.3, 'width': 12, 'height': 0.8},
            'content': {'left': 0.5, 'top': 1.3, 'width': 12, 'height': 5.7}
        }
    }
    
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = presentation_title
    subtitle.text = f"Presented by {presenter_name}"
    apply_template_styles(title, template_choice, is_title=True)
    apply_template_styles(subtitle, template_choice, is_title=False)
    
    def clean_bullet_points(text):
        """Remove dash bullets as PowerPoint will add its own"""
        lines = text.split('\n')
        cleaned_lines = [line[2:].strip() if line.strip().startswith('- ') else line for line in lines]
        return '\n'.join(cleaned_lines)
    
    for i, slide_content in enumerate(slides_content):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        has_image = slide_content.get('image', False)
        if has_image:
            layout_key = 'left_image' if c % 2 == 0 else 'right_image'
            c+=1
        else:
            layout_key = 'no_image'


        layout = layouts[layout_key]
        
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = slide_content['title']
                apply_template_styles(placeholder, template_choice, is_title=True)
                adjust_placeholder_size(placeholder, **layout['title'])
            elif placeholder.placeholder_format.type == 7:  # Content
                cleaned_content = clean_bullet_points(slide_content['content'])
                placeholder.text = cleaned_content
                apply_template_styles(placeholder, template_choice, is_title=False)
                adjust_placeholder_size(placeholder, **layout['content'])
        
        if has_image:
            image_path = get_random_image(image_folder)
            if image_path:
                img_layout = layout['image']
                slide.shapes.add_picture(
                    image_path,
                    left=Inches(img_layout['left']),
                    top=Inches(img_layout['top']),
                    width=Inches(img_layout['width']),
                    height=Inches(img_layout['height'])
                )
    
    delete_first_two_slides(prs)
    output_path = os.path.join('generated_slides', 'generated_presentation.pptx')
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[1])
    prs.save(output_path)
    return output_path

def delete_first_two_slides(presentation):
    slide_ids = [1, 0]
    for slide_id in slide_ids:
        if slide_id < len(presentation.slides):
            xml_slides = presentation.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[slide_id])

create_ppt(slides_content, template_choice, title, authors, 'images')

# minimal
# geometric
# bright_modern
# dark_modern